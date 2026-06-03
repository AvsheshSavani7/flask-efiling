"""
Montana PSC REDDI Docket & Filings Scraper
===========================================
Uses Playwright to authenticate via OKTA SAML SSO, navigate the REDDI portal,
search for dockets, and scrape filing records from the filings table.

Incremental scraping: pass a last_id (e.g. FIL-38222_DOC-69608) to only
process NEW filings/documents since that watermark. Everything before the
matching Filing+Document combination is considered new and gets downloaded.

Install:
    pip install -r requirements.txt
    playwright install chromium

Run:
    python mt_psc_scraper.py --docket 2025.10.078 --case-id DCKT-3556
    python mt_psc_scraper.py --docket 2025.10.078 --case-id DCKT-3556 --last-id FIL-38222_DOC-69608
    python mt_psc_scraper.py --headless --save-json
"""

from __future__ import annotations

import argparse
import base64
import json
import logging
import mimetypes
import os
import sys
import time
from datetime import datetime, timezone
from typing import Any, Dict, List, Optional, Tuple

from dotenv import load_dotenv
from playwright.sync_api import Page, sync_playwright
from playwright.sync_api import TimeoutError as PlaywrightTimeoutError
from tier1_summary_generator import generate_tier1_summary
from aws_utils import build_docket_key, upload_bytes_to_s3

load_dotenv(".env")

# OpenAI fallback for image-only / scanned PDFs (same idea as nm_prc_document_download_extract).
_PDF_OPENAI_OCR_MODEL = os.getenv("MT_PSC_PDF_OCR_MODEL", "gpt-4.1-mini")
_PDF_OPENAI_OCR_MAX_BYTES = int(
    os.getenv("MT_PSC_PDF_OCR_MAX_BYTES", str(15 * 1024 * 1024))
)

LOG_LEVEL = (os.getenv("LOG_LEVEL") or "INFO").upper()
logger = logging.getLogger("mt_psc_scraper")
logger.setLevel(LOG_LEVEL)
if not logger.handlers:
    handler = logging.StreamHandler(sys.stdout)
    handler.setFormatter(
        logging.Formatter("%(asctime)s | %(levelname)s | %(message)s")
    )
    logger.addHandler(handler)
logger.propagate = False

OKTA_SSO_URL = (
    "https://okta.loginmt.com/app/mtgov_pscreddi_1/"
    "exkcjlwf8aBqV2DAD4x7/sso/saml"
)
REDDI_BASE_URL = (
    "https://reddi.mt.gov/prweb/PRAuth/app/reddi/"
    "h6OrH-oHQ9W-Or5woNuiRj_vsZtrQEk-ZZDzaaZeajw*/!STANDARD"
)


def _now_iso() -> str:
    return (
        datetime.now(timezone.utc)
        .replace(microsecond=0)
        .isoformat()
        .replace("+00:00", "Z")
    )


def _normalize_datetime_for_z(value: str) -> str:
    """
    Normalize a scraped date/time string to MM/DD/YYYY.
    """
    if not value:
        return ""

    raw = value.strip()
    if not raw:
        return ""

    formats = (
        "%m/%d/%Y",
        "%m/%d/%Y %I:%M %p",
        "%m/%d/%Y %H:%M",
        "%Y-%m-%d",
        "%Y-%m-%d %H:%M:%S",
    )
    for fmt in formats:
        try:
            parsed = datetime.strptime(raw, fmt)
            return parsed.strftime("%m/%d/%Y")
        except ValueError:
            continue

    # If unknown format, keep raw value.
    return raw


def _parse_last_id(last_id: Optional[str]) -> Tuple[Optional[str], Optional[str]]:
    """Parse 'FIL-38222_DOC-69608' into (filing_id, doc_id)."""
    if not last_id:
        return None, None
    parts = last_id.split("_", 1)
    filing_id = parts[0] if len(parts) >= 1 else None
    doc_id = parts[1] if len(parts) >= 2 else None
    return filing_id, doc_id


# ---------------------------------------------------------------------------
# OKTA SSO Login
# ---------------------------------------------------------------------------

def _login_via_okta(page: Page, username: str, password: str) -> None:
    """Handle the full OKTA SAML SSO → Montana.gov sign-in flow."""
    logger.info("Navigating to OKTA SSO URL...")
    page.goto(OKTA_SSO_URL, wait_until="domcontentloaded", timeout=60000)
    page.wait_for_timeout(3000)

    current_url = page.url
    logger.info(f"Current URL after OKTA redirect: {current_url}")

    if "reddi.mt.gov" in current_url:
        pega_login = page.locator("text=Login with OKTA").first
        if pega_login.is_visible():
            logger.info(
                "PEGA login page detected, clicking 'Login with OKTA'...")
            pega_login.click()
            page.wait_for_timeout(3000)
            current_url = page.url
            logger.info(f"Redirected to: {current_url}")

    if "okta" in current_url.lower() or "loginmt" in current_url.lower():
        logger.info(
            "Montana.gov OKTA sign-in page detected, entering credentials...")
        _fill_okta_credentials(page, username, password)
    else:
        logger.info("Already authenticated or unexpected page state.")


def _fill_okta_credentials(page: Page, username: str, password: str) -> None:
    """Fill in credentials on the Montana.gov OKTA sign-in form."""
    page.wait_for_timeout(3000)

    username_selectors = [
        "input[name='identifier']",
        "input[autocomplete='username']",
        "input#input27",
        "input[name='username']",
        "input#okta-signin-username",
        "input[type='text']",
    ]

    username_field = None
    for sel in username_selectors:
        try:
            field = page.locator(sel).first
            if field.is_visible(timeout=3000):
                username_field = field
                logger.info(f"Found username field with selector: {sel}")
                break
        except PlaywrightTimeoutError:
            continue

    if not username_field:
        raise RuntimeError("Could not find username field on OKTA login page")

    username_field.fill(username)
    logger.info("Username entered.")

    password_selectors = [
        "input[name='credentials.passcode']",
        "input[autocomplete='current-password']",
        "input#input35",
        "input[name='password']",
        "input#okta-signin-password",
        "input[type='password']",
    ]

    password_field = None
    for sel in password_selectors:
        try:
            field = page.locator(sel).first
            if field.is_visible(timeout=3000):
                password_field = field
                logger.info(f"Found password field with selector: {sel}")
                break
        except PlaywrightTimeoutError:
            continue

    if not password_field:
        raise RuntimeError("Could not find password field on OKTA login page")

    password_field.fill(password)
    logger.info("Password entered.")

    submit_selectors = [
        "input.button-primary[type='submit'][value='Sign in']",
        "input[type='submit'][value='Sign in' i]",
        "input[data-type='save']",
        "button:has-text('Sign in')",
        "input[type='submit']",
        "button[type='submit']",
    ]

    for sel in submit_selectors:
        try:
            btn = page.locator(sel).first
            if btn.is_visible(timeout=3000):
                btn.click()
                logger.info(f"Sign-in button clicked (selector: {sel}).")
                break
        except PlaywrightTimeoutError:
            continue

    logger.info("Waiting for authentication redirect...")
    page.wait_for_timeout(5000)

    for attempt in range(20):
        current_url = page.url
        if "reddi.mt.gov" in current_url:
            logger.info("Successfully redirected to REDDI portal.")
            page.wait_for_timeout(3000)
            return
        logger.info(
            f"  Waiting for REDDI redirect... (attempt {attempt + 1}/20, url: {current_url})")
        page.wait_for_timeout(3000)

    raise RuntimeError(f"Login redirect timed out. Current URL: {page.url}")


# ---------------------------------------------------------------------------
# PEGA Navigation
# ---------------------------------------------------------------------------

def _navigate_to_dockets_filings(page: Page) -> None:
    """Click on 'Dockets & Filings' in the REDDI PEGA sidebar navigation."""
    logger.info("Navigating to Dockets & Filings...")
    page.wait_for_timeout(5000)

    try:
        search_menu = page.locator("a[aria-label='Search']").first
        if search_menu.is_visible(timeout=5000):
            expanded = search_menu.get_attribute("aria-expanded")
            logger.info(f"Search menu aria-expanded={expanded}")
            if expanded == "true":
                logger.info("'Search' menu already expanded, skipping click.")
            else:
                logger.info("Expanding 'Search' sidebar menu...")
                search_menu.click()
                page.wait_for_timeout(2000)
    except PlaywrightTimeoutError:
        logger.info("Could not find Search menu toggle, trying direct click...")

    docket_selectors = [
        "a[aria-label='Dockets & Filings']",
        "a[title='Dockets & Filings']",
        "li[data-test-id='202305101433050388232'] a",
    ]

    clicked = False
    for sel in docket_selectors:
        try:
            elem = page.locator(sel).first
            if elem.is_visible(timeout=5000):
                elem.click()
                logger.info(f"Clicked 'Dockets & Filings' (selector: {sel}).")
                clicked = True
                break
        except PlaywrightTimeoutError:
            continue

    if not clicked:
        logger.info("Trying fallback: scanning all menu-item anchors...")
        anchors = page.locator("a.menu-item-anchor")
        for i in range(anchors.count()):
            anchor = anchors.nth(i)
            title = anchor.get_attribute("title") or ""
            aria = anchor.get_attribute("aria-label") or ""
            if "docket" in title.lower() or "docket" in aria.lower():
                anchor.click()
                logger.info(f"Clicked menu anchor: title='{title}'")
                clicked = True
                break

    if not clicked:
        raise RuntimeError(
            "Could not find 'Dockets & Filings' navigation element.")

    page.wait_for_timeout(3000)


def _get_active_pega_iframe(page: Page, timeout: int = 30000) -> Any:
    """
    Find and return the currently ACTIVE PEGA gadget iframe.
    The active module has style="display: block".
    """
    logger.info("Looking for active PEGA iframe...")

    start = time.time()
    while (time.time() - start) * 1000 < timeout:
        modules = page.locator("div.iframe-wrapper.yui-module")
        for i in range(modules.count()):
            module = modules.nth(i)
            style = module.get_attribute("style") or ""
            if "display: block" in style or "display:block" in style:
                header = module.locator("div.hd").first
                header_title = header.get_attribute(
                    "title") or "" if header.count() else ""
                iframe_el = module.locator("iframe").first
                if iframe_el.count() > 0:
                    frame_name = iframe_el.get_attribute("name") or ""
                    iframe_title = iframe_el.get_attribute("title") or ""
                    logger.info(
                        f"Found active module: header='{header_title}', "
                        f"iframe name={frame_name}, iframe title='{iframe_title}'"
                    )
                    frame = page.frame(name=frame_name)
                    if frame:
                        return frame

        for frame in page.frames:
            if frame == page.main_frame:
                continue
            if frame.url and frame.url != "about:blank":
                try:
                    frame_el = page.locator(
                        f"iframe[name='{frame.name}']").first
                    parent_module = frame_el.locator(
                        "xpath=ancestor::div[contains(@class,'iframe-wrapper')]").first
                    style = parent_module.get_attribute("style") or ""
                    if "display: block" in style or "display:block" in style:
                        logger.info(
                            f"Found active frame via ancestry: name={frame.name}")
                        return frame
                except Exception:
                    pass

        page.wait_for_timeout(2000)

    for frame in page.frames:
        if frame != page.main_frame and frame.url and frame.url != "about:blank":
            logger.warning(
                f"Using fallback frame: name={frame.name}, url={frame.url[:80]}")
            return frame

    raise RuntimeError("Could not find any active PEGA iframe")


def _search_docket(page: Page, docket_number: str) -> Any:
    """Enter the docket number and submit. Returns the iframe Frame."""
    logger.info(f"Searching for docket: {docket_number}")

    frame = _get_active_pega_iframe(page)
    frame.wait_for_timeout(3000)

    docket_input_selectors = [
        "input[name*='DocketNumber']",
        "input[placeholder*='XXXX.XX.XXX']",
        "input[placeholder*='docket number' i]",
    ]

    docket_input = None
    for sel in docket_input_selectors:
        try:
            field = frame.locator(sel).first
            if field.is_visible(timeout=5000):
                docket_input = field
                logger.info(f"Found docket input with selector: {sel}")
                break
        except PlaywrightTimeoutError:
            continue

    if not docket_input:
        raise RuntimeError(
            "Could not find docket number input field in iframe.")

    # Use click + type (not fill) to simulate real user input.
    # PEGA's JS event handlers (data-change, data-keydown) only fire
    # on actual input events, not programmatic value changes.
    docket_input.click()
    frame.wait_for_timeout(500)
    docket_input.fill("")
    docket_input.type(docket_number, delay=50)
    logger.info(f"Docket number '{docket_number}' typed.")
    frame.wait_for_timeout(1000)

    # Trigger PEGA's postValue by dispatching change event, then search
    docket_input.dispatch_event("change")
    frame.wait_for_timeout(1000)

    # Try multiple approaches to trigger the search
    search_triggered = False

    # Approach 1: Click the Search button
    search_btn_selectors = [
        "button.pzbutton:has-text('Search')",
        "button.Strong:has-text('Search')",
        "button:has-text('Search')",
    ]
    for sel in search_btn_selectors:
        try:
            btn = frame.locator(sel).first
            if btn.is_visible(timeout=8000):
                btn.click()
                logger.info(f"Search button clicked (selector: {sel}).")
                search_triggered = True
                break
        except PlaywrightTimeoutError:
            continue

    # Approach 2: Press Enter on the input
    if not search_triggered:
        logger.info("Search button not found, pressing Enter...")
        docket_input.press("Enter")
        logger.info("Enter key pressed on docket input.")
        search_triggered = True

    # Wait for results to load
    frame.wait_for_timeout(5000)

    # Poll for search results (DCKT or FIL rows in the results table)
    for attempt in range(6):
        try:
            result_row = frame.locator(
                "tr[oaargs*='DCKT-'], tr[oaargs*='FIL-']").first
            if result_row.is_visible(timeout=5000):
                logger.info("Search results loaded.")
                return frame
        except PlaywrightTimeoutError:
            pass
        logger.info(
            f"  Waiting for search results (attempt {attempt + 1}/6)...")
        frame.wait_for_timeout(3000)

    # If still no results, return frame anyway and let _click_case handle it
    logger.warning("Search results may not have loaded, proceeding anyway...")
    return frame


def _click_case(frame, case_id: str) -> None:
    """Click on a specific REDDI Case ID in the search results table."""
    logger.info(f"Looking for case: {case_id}")

    selectors = [
        f"a:has-text('{case_id}')",
        f"td a:has-text('{case_id}')",
        f"a:text-is('{case_id}')",
    ]

    # Retry up to 3 times with increasing waits (headless can be slow)
    for attempt in range(3):
        for sel in selectors:
            try:
                link = frame.locator(sel).first
                if link.is_visible(timeout=10000):
                    link.click()
                    logger.info(
                        f"Clicked on case {case_id} (attempt {attempt + 1}).")
                    frame.wait_for_timeout(4000)
                    return
            except PlaywrightTimeoutError:
                continue

        logger.info(
            f"  Case {case_id} not found yet (attempt {attempt + 1}/3), waiting...")
        frame.wait_for_timeout(5000)

    raise RuntimeError(
        f"Could not find case link for {case_id} in search results.")


def _click_filings_sidebar(page: Page, case_id: str) -> Any:
    """Click 'Filings' in the case detail sidebar. Returns the case iframe Frame."""
    logger.info("Looking for case detail iframe...")
    page.wait_for_timeout(5000)

    case_frame = _get_active_pega_iframe(page, timeout=20000)
    logger.info(f"Found active case iframe: {case_frame.name}")

    logger.info("Clicking 'Filings' in sidebar...")
    filing_selectors = [
        "a:has-text('Filings')",
        "button:has-text('Filings')",
        "text=Filings",
    ]

    for sel in filing_selectors:
        try:
            elems = case_frame.locator(sel)
            for i in range(elems.count()):
                elem = elems.nth(i)
                text = (elem.text_content() or "").strip()
                if text == "Filings" or text.startswith("Filings"):
                    elem.click()
                    logger.info(
                        "Clicked 'Filings' sidebar item in case iframe.")
                    case_frame.wait_for_timeout(4000)
                    return case_frame
        except PlaywrightTimeoutError:
            continue

    for sel in filing_selectors:
        try:
            elems = page.locator(sel)
            for i in range(elems.count()):
                elem = elems.nth(i)
                text = (elem.text_content() or "").strip()
                if text == "Filings" or text.startswith("Filings"):
                    elem.click()
                    logger.info("Clicked 'Filings' in main page.")
                    page.wait_for_timeout(4000)
                    return case_frame
        except PlaywrightTimeoutError:
            continue

    raise RuntimeError("Could not find 'Filings' in the case sidebar.")


def _extract_case_details(frame) -> Dict[str, str]:
    """Extract case-level details from the case detail page (inside iframe)."""
    details = {}
    detail_fields = [
        "Status", "Docket Number", "Filing Type", "Docket Type",
        "Filing On Behalf Of", "Created", "Updated"
    ]
    for field_name in detail_fields:
        try:
            label = frame.locator(f"text='{field_name}'").first
            if label.is_visible(timeout=1000):
                parent = label.locator("..")
                value_el = parent.locator("span, div, td").last
                val = (value_el.text_content() or "").strip()
                if val and val != field_name:
                    details[field_name] = val
        except (PlaywrightTimeoutError, Exception):
            continue
    return details


# ---------------------------------------------------------------------------
# Filing & Document Extraction (with watermark stop)
# ---------------------------------------------------------------------------

def _scrape_filings_table(
    frame, stop_filing_id: Optional[str], stop_doc_id: Optional[str], row_number: Optional[int], docket_number: Optional[str]
) -> Tuple[List[Dict[str, Any]], bool]:
    """
    Scrape filings from the PEGA grid. For each filing, expand to get documents.
    Stop when we hit the watermark (stop_filing_id + stop_doc_id).

    Returns:
        (filings_list, reached_watermark)
    """
    frame.wait_for_timeout(3000)

    all_filing_rows = frame.evaluate(r"""() => {
        const results = [];
        const table = document.querySelector('table[pl_prop="D_ChildFILCases.pxResults"]');
        if (!table) return results;

        const rows = table.querySelectorAll('tbody > tr[oaargs*="FIL-"]');
        for (const row of rows) {
            if (row.querySelector('table')) continue;
            const cells = row.querySelectorAll(':scope > td');
            if (cells.length < 7) continue;

            const caseIdLink = cells[1].querySelector('a');
            results.push({
                case_id: caseIdLink ? caseIdLink.textContent.trim() : cells[1].textContent.trim(),
                document_indexes: cells[2].textContent.trim(),
                filing_type: cells[3].textContent.trim(),
                description: cells[4].textContent.trim(),
                filed_by: cells[5].textContent.trim(),
                created_date: cells[6].textContent.trim(),
            });
        }
        return results;
    }""")

    logger.info(f"Found {len(all_filing_rows)} total filings in PEGA grid.")

    if stop_filing_id:
        logger.info(f"Watermark: stop at {stop_filing_id}_{stop_doc_id}")

    new_filings = []
    reached_watermark = False

    for idx, filing in enumerate(all_filing_rows):
        filing["row_number"] = row_number
        filing["docket_number"] = docket_number
        fid = filing.get("case_id", "")

        # If this filing IS the watermark filing, we need to check documents
        if stop_filing_id and fid == stop_filing_id:
            docs = _expand_and_extract_documents(frame, idx)
            new_docs = []
            for doc in docs:
                if stop_doc_id and doc.get("document_id", "") == stop_doc_id:
                    reached_watermark = True
                    break
                new_docs.append(doc)

            if new_docs:
                filing["documents"] = new_docs
                new_filings.append(filing)

            reached_watermark = True
            logger.info(
                f"  Reached watermark at {fid}, kept {len(new_docs)} new docs from this filing.")
            break

        # If this filing is BEFORE the watermark, it's new — extract all its docs
        if stop_filing_id and fid != stop_filing_id:
            docs = _expand_and_extract_documents(frame, idx)
            filing["documents"] = docs
            new_filings.append(filing)
            logger.info(f"  New filing {fid}: {len(docs)} documents")
            continue

        # No watermark set — extract everything
        docs = _expand_and_extract_documents(frame, idx)
        filing["documents"] = docs
        new_filings.append(filing)
        if docs:
            logger.info(f"  {fid}: {len(docs)} documents")

    if stop_filing_id and not reached_watermark:
        logger.warning(
            f"Watermark {stop_filing_id} not found in filings — extracted all.")

    logger.info(f"New filings to process: {len(new_filings)}")
    return new_filings, reached_watermark


def _expand_and_extract_documents(frame, row_index: int) -> List[Dict[str, Any]]:
    """
    Click the expand button on a filing row, extract the documents
    sub-table, then collapse.
    """
    documents = []
    try:
        expand_btns = frame.locator(
            "table[pl_prop='D_ChildFILCases.pxResults'] "
            "tbody > tr[oaargs*='FIL-'] > td.expandPane span[data-ctl='expCollIcon']"
        )
        btn_count = expand_btns.count()
        if row_index >= btn_count:
            return documents

        btn = expand_btns.nth(row_index)
        btn_class = btn.get_attribute("class") or ""

        if "expandRowDetails" in btn_class:
            btn.click()
            frame.wait_for_timeout(2000)

        documents = frame.evaluate(r"""(rowIdx) => {
            const results = [];
            const docTables = document.querySelectorAll(
                'table[pl_prop="D_CaseDocuments.pxResults"]'
            );
            if (docTables.length === 0) return results;

            const docTable = docTables[docTables.length - 1];
            const rows = docTable.querySelectorAll('tbody > tr[oaargs]');

            for (const row of rows) {
                const cells = row.querySelectorAll(':scope > td');
                if (cells.length < 5) continue;

                const idLink = cells[0].querySelector('a');
                results.push({
                    document_id: idLink ? idLink.textContent.trim() : cells[0].textContent.trim(),
                    index: cells[1].textContent.trim(),
                    type: cells[2].textContent.trim(),
                    name: cells[3].textContent.trim(),
                    filed_on: cells[4].textContent.trim(),
                });
            }
            return results;
        }""", row_index)

        try:
            btn = expand_btns.nth(row_index)
            btn_class = btn.get_attribute("class") or ""
            if "collapseRowDetails" in btn_class:
                btn.click()
                frame.wait_for_timeout(500)
        except Exception:
            pass

    except Exception as e:
        logger.warning(f"Could not expand row {row_index} for documents: {e}")

    return documents


# ---------------------------------------------------------------------------
# Flatten filings → document-level records
# ---------------------------------------------------------------------------

def _flatten_filings(filings: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """
    Flatten the nested filings+documents structure into one record per document.
    Filing fields are merged into each document record.
    """
    flat = []
    for filing in filings:
        for doc in filing.get("documents", []):
            flat.append({
                "row_number": filing.get("row_number"),
                "case_id": filing.get("case_id", ""),
                "document_indexes": filing.get("document_indexes", ""),
                "filing_type": filing.get("filing_type", ""),
                "description": filing.get("description", ""),
                "filed_by": filing.get("filed_by", ""),
                "created_date": _normalize_datetime_for_z(filing.get("created_date", "")),
                "document_id": doc.get("document_id", ""),
                "document_index": doc.get("index", ""),
                "document_type": doc.get("type", ""),
                "document_name": doc.get("name", ""),
                "document_filed_on": _normalize_datetime_for_z(
                    doc.get("filed_on", "")
                ),
                "document_filename": doc.get("filename", ""),
                "extracted_text": doc.get("extracted_text", ""),
                "docket_number": filing.get("docket_number", ""),
                "s3_key": doc.get("s3_key", ""),
                "s3_url": doc.get("s3_url", ""),
            })
    return flat


# ---------------------------------------------------------------------------
# Text Extraction (PDF, DOCX, XLSX, PPTX, ZIP)
# ---------------------------------------------------------------------------

def _extract_pdf_text_pypdf2(file_path: str) -> str:
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(file_path)
        parts = []
        for page in reader.pages:
            try:
                text = page.extract_text()
                if text:
                    parts.append(text)
            except Exception:
                continue
        return "\n".join(parts)
    except Exception as e:
        logger.warning(f"PyPDF2 extraction failed for {file_path}: {e}")
        return ""


def _extract_pdf_text_pymupdf(file_path: str) -> str:
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return ""
    try:
        doc = fitz.open(file_path)
        parts: List[str] = []
        try:
            for page in doc:
                t = page.get_text()
                if t and t.strip():
                    parts.append(t)
        finally:
            doc.close()
        return "\n".join(parts)
    except Exception as e:
        logger.debug("PyMuPDF extraction failed for %s: %s", file_path, e)
        return ""


def _extract_pdf_text_openai_ocr(file_path: str) -> str:
    api_key = os.getenv("OPENAI_API_KEY_DOCKET")
    if not api_key:
        return ""
    try:
        with open(file_path, "rb") as f:
            pdf_bytes = f.read()
    except OSError as e:
        logger.debug("Could not read PDF for OpenAI OCR %s: %s", file_path, e)
        return ""
    if not pdf_bytes.startswith(b"%PDF"):
        return ""
    if len(pdf_bytes) > _PDF_OPENAI_OCR_MAX_BYTES:
        logger.info(
            "  OpenAI PDF OCR skipped (file too large): %s bytes > %s",
            len(pdf_bytes),
            _PDF_OPENAI_OCR_MAX_BYTES,
        )
        return ""
    doc_name = os.path.basename(file_path)
    try:
        from openai import OpenAI

        client = OpenAI(api_key=api_key)
        pdf_b64 = base64.b64encode(pdf_bytes).decode("utf-8")
        response = client.responses.create(
            model=_PDF_OPENAI_OCR_MODEL,
            input=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "input_text",
                            "text": (
                                "Extract all readable text from this PDF. "
                                "Return only extracted text, preserving order as best as possible. "
                                "Do not summarize."
                            ),
                        },
                        {
                            "type": "input_file",
                            "filename": doc_name,
                            "file_data": f"data:application/pdf;base64,{pdf_b64}",
                        },
                    ],
                }
            ],
        )
        return (response.output_text or "").strip()
    except Exception as e:
        logger.warning("OpenAI PDF OCR failed for %s: %s", doc_name, e)
        return ""


def _extract_text_from_pdf(file_path: str) -> str:
    short = os.path.basename(file_path)
    text = _extract_pdf_text_pypdf2(file_path)
    if (text or "").strip():
        return text
    text = _extract_pdf_text_pymupdf(file_path)
    if (text or "").strip():
        logger.info(
            "  PDF text via PyMuPDF fallback: %s (%s chars)", short, len(text)
        )
        return text
    text = _extract_pdf_text_openai_ocr(file_path)
    if (text or "").strip():
        logger.info(
            "  PDF text via OpenAI OCR fallback: %s (%s chars)", short, len(
                text)
        )
        return text
    return ""


def _extract_text_from_docx(file_path: str) -> str:
    try:
        from docx import Document as DocxDocument
        doc = DocxDocument(file_path)
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    parts.append("\t".join(cells))
        return "\n".join(parts)
    except Exception as e:
        logger.warning(f"DOCX extraction failed for {file_path}: {e}")
        return ""


def _extract_text_from_xlsx(file_path: str) -> str:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True, data_only=True)
        parts = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    parts.append("\t".join(cells))
        wb.close()
        return "\n".join(parts)
    except Exception as e:
        logger.warning(f"XLSX extraction failed for {file_path}: {e}")
        return ""


def _extract_text_from_pptx(file_path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip()
                                 for c in row.cells if c.text.strip()]
                        if cells:
                            parts.append("\t".join(cells))
        return "\n".join(parts)
    except Exception as e:
        logger.warning(f"PPTX extraction failed for {file_path}: {e}")
        return ""


def _extract_text_from_zip(zip_path: str) -> str:
    """Extract text from all supported files inside a zip archive."""
    import tempfile
    import zipfile

    extractors = {
        ".pdf": _extract_text_from_pdf,
        ".docx": _extract_text_from_docx,
        ".xlsx": _extract_text_from_xlsx,
        ".pptx": _extract_text_from_pptx,
    }
    combined = []
    try:
        with zipfile.ZipFile(zip_path, "r") as zf:
            supported = [
                n for n in zf.namelist()
                if os.path.splitext(n)[1].lower() in extractors
            ]
            if not supported:
                logger.info(
                    f"  No supported files inside {os.path.basename(zip_path)}")
                return ""
            with tempfile.TemporaryDirectory() as tmpdir:
                for name in sorted(supported):
                    ext = os.path.splitext(name)[1].lower()
                    extracted_path = zf.extract(name, tmpdir)
                    text = extractors[ext](extracted_path)
                    if text:
                        combined.append(text)
    except Exception as e:
        logger.warning(f"ZIP extraction failed for {zip_path}: {e}")
    return "\n\n".join(combined)


def _extract_text(file_path: str) -> str:
    """Route to the correct extractor based on file extension."""
    ext = os.path.splitext(file_path)[1].lower()
    extractors = {
        ".pdf": _extract_text_from_pdf,
        ".docx": _extract_text_from_docx,
        ".xlsx": _extract_text_from_xlsx,
        ".pptx": _extract_text_from_pptx,
        ".zip": _extract_text_from_zip,
    }
    extractor = extractors.get(ext)
    if extractor:
        return extractor(file_path)
    logger.info(
        f"  Unsupported file type ({ext}): {os.path.basename(file_path)}")
    return ""


# ---------------------------------------------------------------------------
# Document Download
# ---------------------------------------------------------------------------

def _download_all_documents(
    page: Page, case_frame, filings: List[Dict[str, Any]], download_dir: str
) -> None:
    """
    For each document in each filing, click the doc link to open its detail,
    then click 'Download Document' to save the file.
    """
    total_docs = sum(len(f.get("documents", [])) for f in filings)
    logger.info(f"Downloading {total_docs} documents to {download_dir}...")
    downloaded = 0

    for filing in filings:
        for doc in filing.get("documents", []):
            doc_id = doc.get("document_id", "")
            if not doc_id:
                continue

            try:
                doc_link = case_frame.locator(f"a:has-text('{doc_id}')").first
                try:
                    if not doc_link.is_visible(timeout=3000):
                        _expand_filing_by_id(
                            case_frame, filing.get("case_id", ""))
                        doc_link = case_frame.locator(
                            f"a:has-text('{doc_id}')").first
                        if not doc_link.is_visible(timeout=3000):
                            logger.warning(
                                f"  {doc_id}: link not visible, skipping")
                            continue
                except PlaywrightTimeoutError:
                    logger.warning(f"  {doc_id}: link not found, skipping")
                    continue

                doc_link.click()
                page.wait_for_timeout(5000)

                doc_frame = _get_active_pega_iframe(page, timeout=20000)

                # Wait for document detail page to fully render
                download_btn = doc_frame.locator(
                    "button:has-text('Download Document')")
                try:
                    # Poll for the button to appear (PEGA SPA can be slow)
                    download_btn.wait_for(state="visible", timeout=15000)
                except PlaywrightTimeoutError:
                    # Retry: the iframe may have changed
                    logger.info(f"  {doc_id}: retrying iframe lookup...")
                    page.wait_for_timeout(3000)
                    doc_frame = _get_active_pega_iframe(page, timeout=10000)
                    download_btn = doc_frame.locator(
                        "button:has-text('Download Document')")

                try:
                    if download_btn.is_visible(timeout=5000):
                        with page.expect_download(timeout=30000) as download_info:
                            download_btn.click()
                        dl = download_info.value

                        suggested = dl.suggested_filename or f"{doc_id}.pdf"
                        safe_name = f"{doc_id}_{suggested}"
                        save_path = os.path.join(download_dir, safe_name)
                        dl.save_as(save_path)

                        doc["filename"] = suggested
                        downloaded += 1
                        logger.info(f"  {doc_id}: downloaded → {safe_name}")

                        extracted = _extract_text(save_path)
                        doc["extracted_text"] = extracted
                        if extracted:
                            logger.info(
                                f"  {doc_id}: extracted {len(extracted)} chars")
                        else:
                            logger.info(f"  {doc_id}: no text extracted")

                        try:
                            with open(save_path, "rb") as f:
                                file_bytes = f.read()
                            content_type = mimetypes.guess_type(
                                save_path)[0] or "application/octet-stream"
                            s3_key = build_docket_key(safe_name)
                            s3_result = upload_bytes_to_s3(
                                file_bytes, s3_key, content_type=content_type)
                            doc["s3_key"] = s3_result["key"]
                            doc["s3_url"] = s3_result["url"]
                            logger.info(
                                f"  {doc_id}: uploaded to S3 → {s3_result['url']}")
                        except Exception as s3_exc:
                            doc["s3_upload_error"] = str(s3_exc)
                            logger.warning(
                                f"  {doc_id}: S3 upload failed: {s3_exc}")
                    else:
                        logger.warning(
                            f"  {doc_id}: Download button not visible")
                except PlaywrightTimeoutError:
                    logger.warning(f"  {doc_id}: download timed out")
                except Exception as e:
                    logger.warning(f"  {doc_id}: download error: {e}")

                close_btn = doc_frame.locator("button:has-text('Close')").first
                try:
                    if close_btn.is_visible(timeout=2000):
                        close_btn.click()
                        page.wait_for_timeout(2000)
                except Exception:
                    pass

            except Exception as e:
                logger.warning(f"  {doc_id}: error: {e}")
                continue

    logger.info(f"Downloaded {downloaded}/{total_docs} documents.")


def _expand_filing_by_id(frame, filing_case_id: str) -> None:
    """Expand a specific filing row by its Case ID to reveal its documents."""
    try:
        row = frame.locator(f"tr[oaargs*='{filing_case_id}']").first
        expand_btn = row.locator("span[data-ctl='expCollIcon']").first
        btn_class = expand_btn.get_attribute("class") or ""
        if "expandRowDetails" in btn_class:
            expand_btn.click()
            frame.wait_for_timeout(2000)
    except Exception:
        pass


# ---------------------------------------------------------------------------
# Main Entry Point
# ---------------------------------------------------------------------------

def scrape_mt_psc(
    docket_number: str = "2025.10.078",
    case_id: str = "DCKT-3556",
    last_id: Optional[str] = None,
    username: Optional[str] = None,
    password: Optional[str] = None,
    headless: bool = True,
    save_json: bool = False,
    row_number: Optional[int] = None,
) -> Dict[str, Any]:
    """
    Main scraper entry point.

    Args:
        docket_number: Docket number to search (format: YYYY.NN.NNN)
        case_id: REDDI Case ID to open (e.g. DCKT-3556)
        last_id: Watermark ID (e.g. FIL-38222_DOC-69608). Only filings/docs
                 BEFORE this combination are considered new and downloaded.
                 If None, all filings are processed.
        username: OKTA username (falls back to env MT_PSC_USERNAME)
        password: OKTA password (falls back to env MT_PSC_PASSWORD)
        headless: Run browser in headless mode
        save_json: Also save results to a JSON file

    Returns:
        Dict with scrape results
    """
    username = username or os.getenv("MT_PSC_USERNAME", "")
    password = password or os.getenv("MT_PSC_PASSWORD", "")

    if not username or not password:
        raise ValueError(
            "Username and password required. Set MT_PSC_USERNAME/MT_PSC_PASSWORD "
            "env vars or pass as arguments."
        )

    stop_filing_id, stop_doc_id = _parse_last_id(last_id)

    flat_records = []
    result = {}

    download_dir = os.path.join(
        os.getcwd(), "mt_psc_downloads", docket_number.replace(".", "_"))
    os.makedirs(download_dir, exist_ok=True)
    logger.info(f"Download directory: {download_dir}")

    with sync_playwright() as p:
        launch_args = ["--disable-blink-features=AutomationControlled"]
        if headless:
            launch_args.append("--headless=new")
        browser = p.chromium.launch(headless=headless, args=launch_args)
        context = browser.new_context(
            viewport={"width": 1280, "height": 800},
            user_agent=(
                "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
                "(KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
            ),
            accept_downloads=True,
        )
        page = context.new_page()

        try:
            # Step 1: Login via OKTA SSO
            _login_via_okta(page, username, password)

            # Step 2: Navigate to Dockets & Filings
            _navigate_to_dockets_filings(page)

            # Step 3: Search by docket number
            search_frame = _search_docket(page, docket_number)

            # Step 4: Click on the specific case
            _click_case(search_frame, case_id)

            # Step 5: Click Filings sidebar
            case_frame = _click_filings_sidebar(page, case_id)

            # Step 6: Extract case-level details
            case_details = _extract_case_details(case_frame)
            result["case_details"] = case_details
            logger.info(f"Case details: {json.dumps(case_details, indent=2)}")

            # Step 7: Scrape filings + documents (stop at watermark)
            filings, reached_watermark = _scrape_filings_table(
                case_frame, stop_filing_id, stop_doc_id, row_number, docket_number
            )
            result["total_filings"] = len(filings)
            result["reached_watermark"] = reached_watermark

            # Step 8: Download new documents only
            if filings:
                _download_all_documents(
                    page, case_frame, filings, download_dir)
            else:
                logger.info("No new filings to download.")

            # Flatten: one record per document with filing fields merged in
            flat_records = _flatten_filings(filings)

            result["filings"] = flat_records
            result["download_dir"] = download_dir
            result["success"] = True

            if flat_records:
                new_watermark = (
                    f"{flat_records[0]['case_id']}_{flat_records[0]['document_id']}"
                )
                result["new_last_id"] = new_watermark
                logger.info(f"New watermark: {new_watermark}")

            logger.info(
                f"Scraped {len(filings)} new filings, "
                f"{len(flat_records)} documents for docket {docket_number}")

            if save_json:
                out_file = f"mt_psc_{docket_number.replace('.', '_')}.json"
                with open(out_file, "w", encoding="utf-8") as f:
                    json.dump(flat_records, f, indent=2, ensure_ascii=False)
                logger.info(f"Results saved to {out_file}")
                result["json_file"] = out_file

        except Exception as e:
            logger.error(f"Scraper error: {e}")
            result["error"] = str(e)
            import traceback
            traceback.print_exc()
        finally:
            context.close()
            browser.close()
    print(f"Flat records: {flat_records}")

    return flat_records[::-1]


def main():
    parser = argparse.ArgumentParser(
        description="Montana PSC REDDI Docket & Filings Scraper"
    )
    parser.add_argument(
        "--docket", default="2025.10.078",
        help="Docket number to search (default: 2025.10.078)",
    )
    parser.add_argument(
        "--case-id", default="DCKT-3556",
        help="REDDI Case ID to open (default: DCKT-3556)",
    )
    parser.add_argument(
        "--last-id", default="FIL-38392_DOC-69839",
        help="Watermark: FIL-xxxxx_DOC-xxxxx. Only process filings BEFORE this.",
    )
    parser.add_argument(
        "--username", default=None,
        help="OKTA username (or set MT_PSC_USERNAME env var)",
    )
    parser.add_argument(
        "--password", default=None,
        help="OKTA password (or set MT_PSC_PASSWORD env var)",
    )
    parser.add_argument(
        "--headless", action="store_true", default=False,
        help="Run in headless mode",
    )
    parser.add_argument(
        "--save-json", action="store_true", default=False,
        help="Save results to a JSON file",
    )
    args = parser.parse_args()

    records = scrape_mt_psc(
        docket_number=args.docket,
        case_id=args.case_id,
        last_id=args.last_id,
        username=args.username,
        password=args.password,
        headless=args.headless,
        save_json=args.save_json,
    )

    if records:
        print(f"\nSuccess! Scraped {len(records)} documents.")
        for rec in records:
            print(f"  - {rec.get('case_id')} | {rec.get('document_id')} | "
                  f"{rec.get('filing_type')} | {rec.get('description')}")

        print("\nGenerating tier1 summaries and saving to MongoDB...")
        for rec in records:
            text = (rec.get("extracted_text") or "").strip()
            if not text:
                print(
                    f"  - {rec.get('document_id')}: skipped (no extracted_text)")
                continue

            metadata = {
                "document_id": rec.get("document_id", ""),
                "date": rec.get("document_filed_on") or rec.get("created_date") or "",
                "document_type": rec.get("document_type") or rec.get("filing_type") or "N/A",
                "additional_info": rec.get("description", ""),
                "on_behalf_of": rec.get("filed_by", ""),
                "docket_number": rec.get("docket_number") or args.docket,
                "docket_type": "mt-psc",
            }
            result = generate_tier1_summary(metadata=metadata, text=text)
            status = result.get("status", "unknown")
            if result.get("error"):
                print(
                    f"  - {rec.get('document_id')}: error - {result.get('error')}")
            else:
                print(
                    f"  - {rec.get('document_id')}: {status} (summary_length={result.get('summary_length', 0)})")
    else:
        print("\nNo new documents found (or scraper failed).")
        sys.exit(1)


if __name__ == "__main__":
    main()
